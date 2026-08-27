import os
import re
import pandas as pd
import streamlit as st
from pypdf import PdfReader
import docx
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# Page Setup - 
st.set_page_config(
    page_title="Apna Private Search Engine", 
    page_icon="⚡", 
    layout="wide"
)

DOCS_DIR = "documents"
os.makedirs(DOCS_DIR, exist_ok=True)

def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text_chunks = []

    try:
        if ext in ['.txt', '.md', '.py', '.c', '.cpp', '.java', '.html', '.css', '.js', '.json']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                if content.strip():
                    text_chunks.append({"page": 1, "text": content})

        elif ext == '.pdf':
            reader = PdfReader(file_path)
            for page_num, page in enumerate(reader.pages, start=1):
                t = page.extract_text()
                if t and t.strip():
                    text_chunks.append({"page": page_num, "text": t})

        elif ext == '.docx':
            doc = docx.Document(file_path)
            full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if full_text:
                text_chunks.append({"page": 1, "text": full_text})

        elif ext == '.pptx':
            prs = Presentation(file_path)
            for idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text_chunks.append({"page": idx, "text": "\n".join(slide_text)})
    except Exception as e:
        st.warning(f"Error reading {os.path.basename(file_path)}: {e}")

    return text_chunks

@st.cache_data
def build_index(folder_path):
    chunk_records = []
    
    if not os.path.exists(folder_path):
        return None, None, None

    files = [f for f in os.listdir(folder_path) if not f.startswith('.')]

    for file_name in files:
        file_path = os.path.join(folder_path, file_name)
        chunks = extract_text(file_path)
        for c in chunks:
            chunk_records.append({
                "file_name": file_name,
                "file_path": file_path,
                "file_type": os.path.splitext(file_name)[1].upper().replace('.', ''),
                "page_no": c["page"],
                "text": c["text"]
            })

    if not chunk_records:
        return None, None, None

    df = pd.DataFrame(chunk_records)
    vectorizer = TfidfVectorizer(stop_words="english")
    tfidf_matrix = vectorizer.fit_transform(df["text"])

    return df, vectorizer, tfidf_matrix

def highlight_words(text, query):
    clean_query = query.replace('"', '').strip()
    keywords = [re.escape(word) for word in clean_query.split() if len(word) > 1]
    if not keywords:
        return text
    pattern = re.compile(r'(' + '|'.join(keywords) + r')', re.IGNORECASE)
    return pattern.sub(r'<mark style="background-color: #fde047; color: #000; padding: 2px 4px; border-radius: 3px;">\1</mark>', text)

# Sidebar UI
st.sidebar.title("🛠️ Control Panel")

st.sidebar.subheader("📤 Upload New Files")
uploaded_files = st.sidebar.file_uploader(
    "Choose PDF, Word, PPT, or Code files", 
    type=['pdf', 'docx', 'pptx', 'txt', 'py', 'java', 'c', 'cpp', 'html', 'js'],
    accept_multiple_files=True
)

if uploaded_files:
    for u_file in uploaded_files:
        save_path = os.path.join(DOCS_DIR, u_file.name)
        with open(save_path, "wb") as f:
            f.write(u_file.getbuffer())
    st.sidebar.success(f"Successfully uploaded {len(uploaded_files)} file(s)!")
    st.cache_data.clear()

st.sidebar.markdown("---")

# Modern Gradient Title Header
st.markdown("""
    <h1 style='text-align: left; font-size: 2.8rem; font-weight: 800; background: linear-gradient(90deg, #3B82F6 0%, #8B5CF6 50%, #EC4899 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent;'>
        ⚡ Apna Private Search Engine
    </h1>
""", unsafe_allow_html=True)

st.caption("Next-Gen Multi-Format Search across PDFs, Word, PPT, & Code Files")

df, vectorizer, tfidf_matrix = build_index(DOCS_DIR)

if df is None or df.empty:
    st.warning("The `documents/` folder is empty. Please upload files using the sidebar.")
else:
    file_types = ["ALL"] + list(df['file_type'].unique())
    selected_type = st.sidebar.selectbox("📁 Filter by File Type:", file_types)
    st.sidebar.metric("Total Document Chunks", len(df))

    # Search Bar & Button Grid System
    col1, col2 = st.columns([5, 1])
    
    with col1:
        query = st.text_input("Search Query", placeholder="Enter topic e.g., Ternary Operators, array, function", label_visibility="collapsed")
    
    with col2:
        search_btn = st.button("🔍 Search", use_container_width=True)

    # Search Trigger (Enter Press OR Button Click)
    if query or search_btn:
        clean_query = query.replace('"', '').strip()
        
        if clean_query:
            query_vec = vectorizer.transform([clean_query])
            similarity = cosine_similarity(query_vec, tfidf_matrix).flatten()

            df["score"] = similarity
            results_df = df[df["score"] > 0.05].copy()

            # Strict Multi-word matching
            if len(clean_query.split()) > 1:
                words = clean_query.lower().split()
                results_df = results_df[results_df['text'].apply(lambda t: all(w in t.lower() for w in words))]

            results_df = results_df.sort_values(by="score", ascending=False)

            if selected_type != "ALL":
                results_df = results_df[results_df['file_type'] == selected_type]

            st.subheader(f"Matching Results ({len(results_df)} matches found)")

            if results_df.empty:
                st.info("No relevant matches found for your query.")
            else:
                for _, row in results_df.head(15).iterrows():
                    score_pct = round(float(row['score']) * 100, 1)
                    file_badge = f"[{row['file_type']}]"
                    highlighted_text = highlight_words(row["text"].strip()[:1500], query)
                    
                    with st.expander(f"📄 {file_badge} {row['file_name']} (Page/Slide {row['page_no']}) — Relevance: {score_pct}%"):
                        st.markdown(highlighted_text, unsafe_allow_html=True)
                        st.markdown("---")
                        
                        if os.path.exists(row['file_path']):
                            with open(row['file_path'], "rb") as file_bytes:
                                st.download_button(
                                    label=f"📥 Download {row['file_name']}",
                                    data=file_bytes,
                                    file_name=row['file_name'],
                                    mime="application/octet-stream",
                                    key=f"dl_{row['file_name']}_{row['page_no']}"
                                )